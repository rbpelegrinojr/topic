from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy


DARK_BLUE   = RGBColor(0x0D, 0x1B, 0x2A)   # background / title bg
MID_BLUE    = RGBColor(0x1B, 0x48, 0x8B)   # accent bar
ACCENT_BLUE = RGBColor(0x26, 0x8B, 0xD2)   # highlights
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xEC, 0xEF, 0xF4)
YELLOW      = RGBColor(0xFF, 0xD1, 0x66)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
RED_SOFT    = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


def add_rect(slide, left, top, width, height, fill_color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  base_size=14, color=WHITE, bold_first=False):
    """Add multiple paragraphs inside one text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if isinstance(line, tuple):
            run.text, size, bold, col = line
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = col
        else:
            run.text = line
            run.font.size  = Pt(base_size)
            run.font.bold  = bold_first and first
            run.font.color.rgb = color
    return txBox


def slide_chrome(slide, section_tag=""):
    """Draw the common background + top bar + bottom bar."""
    # Full background
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
    # Top accent bar
    add_rect(slide, 0, 0, 13.33, 0.12, ACCENT_BLUE)
    # Bottom bar
    add_rect(slide, 0, 7.25, 13.33, 0.25, MID_BLUE)
    # Section tag bottom-right
    if section_tag:
        add_text_box(slide, section_tag, 10.5, 7.22, 2.7, 0.28,
                     font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def slide_title_bar(slide, title_text, subtitle_text=""):
    """Blue title rectangle near top."""
    add_rect(slide, 0.4, 0.18, 12.53, 0.95, MID_BLUE)
    add_text_box(slide, title_text, 0.55, 0.2, 12.2, 0.9,
                 font_size=28, bold=True, color=WHITE)
    if subtitle_text:
        add_text_box(slide, subtitle_text, 0.55, 1.1, 12.2, 0.4,
                     font_size=15, color=ACCENT_BLUE, bold=False)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 0,  13.33, 0.15, ACCENT_BLUE)
add_rect(slide, 0, 7.35, 13.33, 0.15, ACCENT_BLUE)
add_rect(slide, 1.5, 1.5, 10.33, 4.5, MID_BLUE)
add_rect(slide, 1.55, 1.55, 10.23, 0.08, ACCENT_BLUE)

add_text_box(slide, "Database Normalization Concepts",
             1.7, 1.7, 9.9, 1.2,
             font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "Understanding Normal Forms: 1NF, 2NF, and 3NF",
             1.7, 2.85, 9.9, 0.6,
             font_size=20, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 3.55, 6.33, 0.05, ACCENT_BLUE)

add_text_box(slide, "System Design and Implementation",
             1.7, 3.7, 9.9, 0.5,
             font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, "April 2026",
             1.7, 4.25, 9.9, 0.4,
             font_size=14, color=YELLOW, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "Overview")
slide_title_bar(slide, "What We'll Cover Today")

topics = [
    "  📌  A. What is Normalization?",
    "         •  Definition and Purpose",
    "         •  Why We Normalize",
    "",
    "  📌  B. Normal Forms",
    "         •  First Normal Form (1NF)",
    "         •  Second Normal Form (2NF)",
    "         •  Third Normal Form (3NF)",
    "",
    "  📌  C. Functional Dependencies",
    "         •  Determinants and Dependent Attributes",
    "         •  Identifying Functional Dependencies",
    "",
    "  📌  D. Normalization Process & Best Practices",
]
lines = [(t, 15, "📌" in t, WHITE) for t in topics]
add_multiline(slide, lines, 0.6, 1.3, 12.1, 5.8, base_size=15)

slide.notes_slide.notes_text_frame.text = (
    "Briefly orient students to today's topics. Emphasize that normalization prevents "
    "data anomalies and that each normal form builds on the previous one. Encourage "
    "students to think of each NF as a progressively stricter rule that removes a "
    "specific type of redundancy."
)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 3 — WHAT IS NORMALIZATION?
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "A. What is Normalization?")
slide_title_bar(slide, "What is Normalization?")

add_rect(slide, 0.4, 1.32, 12.5, 0.55, MID_BLUE)
add_text_box(slide,
             "Normalization is the process of organizing a relational database to reduce "
             "data redundancy and improve data integrity.",
             0.55, 1.34, 12.2, 0.52, font_size=14, color=WHITE)

# Left: Why normalize
add_rect(slide, 0.4, 2.0, 5.9, 0.38, RGBColor(0x0D, 0x30, 0x1A))
add_text_box(slide, "✅  Why Normalize?", 0.5, 2.02, 5.7, 0.34,
             font_size=14, bold=True, color=GREEN)
why = [
    "  •  Eliminate duplicate data",
    "  •  Avoid update, insert, and delete anomalies",
    "  •  Ensure data consistency across the database",
    "  •  Simplify queries and maintenance",
    "  •  Enforce relationships through structure",
]
add_multiline(slide, [(t, 13, False, WHITE) for t in why], 0.4, 2.42, 5.9, 2.1)

# Right: Anomalies explained
add_rect(slide, 6.9, 2.0, 5.9, 0.38, RGBColor(0x40, 0x10, 0x10))
add_text_box(slide, "❌  Data Anomalies (Without Normalization)", 7.0, 2.02, 5.7, 0.34,
             font_size=13, bold=True, color=RED_SOFT)
anomalies = [
    ("  Update Anomaly",   "Changing one value requires updating many rows"),
    ("  Insert Anomaly",   "Cannot insert data without unrelated data"),
    ("  Delete Anomaly",   "Deleting a row unintentionally removes other facts"),
]
y_a = 2.42
for title, desc in anomalies:
    add_rect(slide, 6.9, y_a, 5.9, 0.55, RGBColor(0x2A, 0x0D, 0x0D))
    add_text_box(slide, title, 7.0, y_a + 0.02, 5.7, 0.25, font_size=12, bold=True, color=YELLOW)
    add_text_box(slide, desc, 7.0, y_a + 0.27, 5.7, 0.25, font_size=11, color=WHITE)
    y_a += 0.6

# Bottom: UNF example
add_text_box(slide, "Example — Unnormalized (UNF) Table:",
             0.4, 4.62, 12.5, 0.3, font_size=13, bold=True, color=ACCENT_BLUE)
add_rect(slide, 0.4, 4.95, 12.5, 1.05, RGBColor(0x10, 0x25, 0x40))
unf = (
    "  StudentID  |  StudentName  |  Courses\n"
    "  ─────────────────────────────────────────────────────────────\n"
    "    S001      |  Alice        |  Math, Science, English\n"
    "    S002      |  Bob          |  Science, History\n"
    "  ⚠  Courses column contains multiple values — violates atomicity!"
)
add_text_box(slide, unf, 0.55, 4.97, 12.2, 1.0, font_size=12, color=LIGHT_GRAY)

slide.notes_slide.notes_text_frame.text = (
    "Start with the 'why' before the 'how'. Students understand normalization better "
    "when they first see what goes wrong without it. The UNF example shows a multi-valued "
    "column which is the most common first mistake."
)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 4 — FIRST NORMAL FORM (1NF)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "B. Normal Forms")
slide_title_bar(slide, "First Normal Form (1NF) — Atomic Values")

add_rect(slide, 0.4, 1.32, 12.5, 0.72, MID_BLUE)
add_text_box(slide, "Definition", 0.55, 1.34, 3.0, 0.32,
             font_size=14, bold=True, color=YELLOW)
add_text_box(slide,
             "A table is in 1NF if: (1) All attribute values are atomic (indivisible), "
             "(2) There are no repeating groups or arrays, and (3) Each row is uniquely identifiable.",
             0.55, 1.62, 12.2, 0.4, font_size=13, color=WHITE)

# Before / After
add_text_box(slide, "❌  Before 1NF (Violates — repeating groups):",
             0.4, 2.12, 6.1, 0.3, font_size=12, bold=True, color=RED_SOFT)
add_rect(slide, 0.4, 2.44, 6.1, 1.0, RGBColor(0x30, 0x10, 0x10))
before_tbl = (
    "  OrderID  |  CustomerName  |  Items\n"
    "  ──────────────────────────────────────────\n"
    "    1001    |  Alice         |  Laptop, Mouse\n"
    "    1002    |  Bob           |  Keyboard"
)
add_text_box(slide, before_tbl, 0.55, 2.46, 5.9, 0.96, font_size=11, color=LIGHT_GRAY)

add_text_box(slide, "✅  After 1NF (Each cell is atomic):",
             6.85, 2.12, 6.1, 0.3, font_size=12, bold=True, color=GREEN)
add_rect(slide, 6.85, 2.44, 6.1, 1.0, RGBColor(0x0D, 0x30, 0x1A))
after_tbl = (
    "  OrderID  |  CustomerName  |  Item\n"
    "  ─────────────────────────────────────────\n"
    "    1001    |  Alice         |  Laptop\n"
    "    1001    |  Alice         |  Mouse\n"
    "    1002    |  Bob           |  Keyboard"
)
add_text_box(slide, after_tbl, 7.0, 2.46, 5.9, 0.96, font_size=11, color=LIGHT_GRAY)

# Rules summary
add_rect(slide, 0.4, 3.55, 12.5, 0.38, MID_BLUE)
add_text_box(slide, "1NF Rules at a Glance:", 0.55, 3.57, 5.0, 0.34,
             font_size=13, bold=True, color=YELLOW)
rules_1nf = [
    ("  ✅  Each column holds only one value per row (atomic)", 13, False, WHITE),
    ("  ✅  No multi-valued attributes (e.g., comma-separated lists)", 13, False, WHITE),
    ("  ✅  No repeating column groups (e.g., Item1, Item2, Item3…)", 13, False, WHITE),
    ("  ✅  Every row must be uniquely identifiable (Primary Key)", 13, False, WHITE),
]
add_multiline(slide, rules_1nf, 0.4, 3.98, 12.5, 1.6)

add_rect(slide, 0.4, 5.72, 12.5, 0.42, RGBColor(0x14, 0x2D, 0x50))
add_text_box(slide,
             "  💡  Key Insight: 1NF ensures every cell contains a single, "
             "indivisible value — the foundation for all higher normal forms.",
             0.55, 5.74, 12.2, 0.38, font_size=12, bold=True, color=ACCENT_BLUE)

slide.notes_slide.notes_text_frame.text = (
    "The 'Items' column in the before table stores multiple values in one cell — this is the "
    "most common 1NF violation. After applying 1NF, each item gets its own row. This may "
    "seem to create duplication, but that's addressed by higher normal forms."
)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 5 — SECOND NORMAL FORM (2NF)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "B. Normal Forms")
slide_title_bar(slide, "Second Normal Form (2NF) — Eliminating Partial Dependencies")

add_rect(slide, 0.4, 1.32, 12.5, 0.72, MID_BLUE)
add_text_box(slide, "Definition", 0.55, 1.34, 3.0, 0.32,
             font_size=14, bold=True, color=YELLOW)
add_text_box(slide,
             "A table is in 2NF if: (1) It is already in 1NF, AND "
             "(2) Every non-key attribute is FULLY functionally dependent on the ENTIRE primary key. "
             "No partial dependencies allowed.",
             0.55, 1.62, 12.2, 0.4, font_size=13, color=WHITE)

add_text_box(slide,
             "⚠  Partial Dependency: A non-key attribute depends on only PART of a composite key.",
             0.4, 2.12, 12.5, 0.3, font_size=13, bold=True, color=YELLOW)

# Violating table
add_text_box(slide, "❌  Violating Table (Composite PK: OrderID + ProductID):",
             0.4, 2.52, 7.0, 0.3, font_size=12, bold=True, color=RED_SOFT)
add_rect(slide, 0.4, 2.84, 12.5, 0.78, RGBColor(0x10, 0x25, 0x40))
viol_tbl = (
    "  OrderID  |  ProductID  |  ProductName  |  Quantity  |  CustomerName\n"
    "  ─────────────────────────────────────────────────────────────────────────\n"
    "    1001    |    P01      |    Laptop     |     2      |  Alice"
)
add_text_box(slide, viol_tbl, 0.55, 2.86, 12.2, 0.74, font_size=11, color=LIGHT_GRAY)

fds_2nf = [
    ("  {OrderID, ProductID}  →  Quantity         ✅  Full dependency on composite PK", 12, False, WHITE),
    ("  ProductID             →  ProductName      ⚠  Partial — only needs ProductID!", 12, True, YELLOW),
    ("  OrderID               →  CustomerName     ⚠  Partial — only needs OrderID!", 12, True, YELLOW),
]
add_multiline(slide, fds_2nf, 0.4, 3.68, 12.5, 0.96)

# Solution
add_text_box(slide, "✅  Solution — Decompose into 3 tables:",
             0.4, 4.72, 8.0, 0.3, font_size=12, bold=True, color=GREEN)
tables_2nf = [
    ("Orders", "OrderID  |  CustomerName", 0.4, 5.06, 3.9),
    ("Products", "ProductID  |  ProductName", 4.55, 5.06, 3.9),
    ("OrderDetails", "OrderID  |  ProductID  |  Quantity", 8.7, 5.06, 4.2),
]
for tname, cols, lx, ly, lw in tables_2nf:
    add_rect(slide, lx, ly, lw, 0.72, RGBColor(0x0D, 0x30, 0x1A))
    add_text_box(slide, tname, lx + 0.1, ly + 0.02, lw - 0.2, 0.28,
                 font_size=12, bold=True, color=YELLOW)
    add_text_box(slide, cols, lx + 0.1, ly + 0.3, lw - 0.2, 0.36,
                 font_size=11, color=WHITE)

add_rect(slide, 0.4, 5.9, 12.5, 0.38, RGBColor(0x14, 0x2D, 0x50))
add_text_box(slide,
             "  💡  Key Insight: 2NF only applies when the primary key is composite. "
             "If the PK is a single column, a table in 1NF is automatically in 2NF.",
             0.55, 5.92, 12.2, 0.34, font_size=12, bold=True, color=ACCENT_BLUE)

slide.notes_slide.notes_text_frame.text = (
    "Stress that partial dependencies only apply to composite primary keys. "
    "ProductName depends only on ProductID — not the full {OrderID, ProductID} composite key. "
    "The fix is to move it to a separate Products table. 2NF eliminates this data duplication."
)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 6 — THIRD NORMAL FORM (3NF)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "B. Normal Forms")
slide_title_bar(slide, "Third Normal Form (3NF) — Eliminating Transitive Dependencies")

add_rect(slide, 0.4, 1.32, 12.5, 0.72, MID_BLUE)
add_text_box(slide, "Definition", 0.55, 1.34, 3.0, 0.32,
             font_size=14, bold=True, color=YELLOW)
defn = ("A table is in 3NF if: (1) It is already in 2NF, AND "
        "(2) Every non-key attribute is NON-TRANSITIVELY dependent on the primary key.\n"
        "Mnemonic: '...the KEY, the WHOLE KEY, and NOTHING BUT THE KEY.'")
add_text_box(slide, defn, 0.55, 1.62, 12.2, 0.4, font_size=13, color=WHITE)

add_text_box(slide,
             "⚠  Transitive Dependency: If  A → B  and  B → C,  then  A → C  is TRANSITIVE (indirect).",
             0.4, 2.12, 12.5, 0.3, font_size=13, bold=True, color=YELLOW)

# Violating table
add_text_box(slide, "❌  Violating Table (PK: StudentID):",
             0.4, 2.52, 6.0, 0.3, font_size=12, bold=True, color=RED_SOFT)
add_rect(slide, 0.4, 2.84, 12.5, 0.72, RGBColor(0x10, 0x25, 0x40))
tbl = ("  StudentID  |  StudentName  |  DeptID  |  DeptName\n"
       "  ─────────────────────────────────────────────────────────\n"
       "    S001      |  Alice        |   D01    |  Computer Science")
add_text_box(slide, tbl, 0.55, 2.86, 12.2, 0.68, font_size=11, color=LIGHT_GRAY)

fds_3nf = [
    ("  StudentID  →  DeptID          ✅  Direct dependency on PK", 12, False, WHITE),
    ("  DeptID     →  DeptName        ✅  Direct (but DeptID is not the PK!)", 12, False, WHITE),
    ("  StudentID  →  DeptName        ❌  TRANSITIVE — goes through DeptID!", 12, True, YELLOW),
]
add_multiline(slide, fds_3nf, 0.4, 3.64, 8.5, 0.96)

add_rect(slide, 9.05, 3.64, 3.8, 0.96, RGBColor(0x40, 0x10, 0x10))
add_text_box(slide, "Caused Anomalies:", 9.15, 3.66, 3.6, 0.28,
             font_size=12, bold=True, color=RED_SOFT)
add_multiline(slide, [
    ("  ⚠  Update Anomaly", 11, False, WHITE),
    ("  ⚠  Deletion Anomaly", 11, False, WHITE),
    ("  ⚠  Insertion Anomaly", 11, False, WHITE),
], 9.15, 3.96, 3.6, 0.6)

# Solution
add_text_box(slide, "✅  Solution — Decompose into 2 tables:",
             0.4, 4.72, 8.0, 0.3, font_size=12, bold=True, color=GREEN)
add_rect(slide, 0.4, 5.06, 5.8, 0.64, RGBColor(0x0D, 0x30, 0x1A))
add_text_box(slide, "Student Table:", 0.55, 5.08, 5.6, 0.26, font_size=12, bold=True, color=YELLOW)
add_text_box(slide, "StudentID  |  StudentName  |  DeptID",
             0.55, 5.34, 5.6, 0.32, font_size=11, color=WHITE)
add_rect(slide, 7.1, 5.06, 5.8, 0.64, RGBColor(0x0D, 0x30, 0x1A))
add_text_box(slide, "Department Table:", 7.25, 5.08, 5.6, 0.26, font_size=12, bold=True, color=YELLOW)
add_text_box(slide, "DeptID  |  DeptName",
             7.25, 5.34, 5.6, 0.32, font_size=11, color=WHITE)

add_rect(slide, 0.4, 5.84, 12.5, 0.38, RGBColor(0x14, 0x2D, 0x50))
add_text_box(slide,
             "  💡  Key Insight: DeptName now lives in exactly one place. "
             "Changing a department name requires updating only one row.",
             0.55, 5.86, 12.2, 0.34, font_size=12, bold=True, color=ACCENT_BLUE)

slide.notes_slide.notes_text_frame.text = (
    "The 'nothing but the key' mnemonic covers all three normal forms: "
    "1NF=the key, 2NF=the whole key, 3NF=nothing but the key. "
    "Walk through the deletion anomaly — if the last student in a department is deleted, "
    "the department name is lost. After decomposition, DeptName lives in its own table."
)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 7 — FUNCTIONAL DEPENDENCIES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "C. Functional Dependencies")
slide_title_bar(slide, "Functional Dependencies — Concepts and Identification")

add_rect(slide, 0.4, 1.32, 12.5, 0.5, MID_BLUE)
add_text_box(slide,
             "FD  X → Y: For every valid row, if two rows share the same X value, "
             "they MUST have the same Y value.  (\"X functionally determines Y\")",
             0.55, 1.34, 12.2, 0.46, font_size=13, color=WHITE)

# Terms table
terms = [
    ("Determinant",   "Left side of FD — uniquely determines another attr", "StudentID in StudentID → Name"),
    ("Dependent",     "Right side of FD — is determined by another attr",   "Name in StudentID → Name"),
    ("Full FD",       "Dependent relies on the ENTIRE key",                 "Critical for 2NF"),
    ("Partial FD",    "Dependent relies on PART of composite key",          "Violates 2NF"),
    ("Transitive FD", "A → B → C  chain through non-key attribute",         "Violates 3NF"),
]
col_xs_t = [0.4, 2.9, 7.0]
col_ws_t = [2.45, 4.05, 5.85]
add_rect(slide, 0.4, 1.9, 12.9, 0.32, MID_BLUE)
for ci, h in enumerate(["Term", "Definition", "Example"]):
    cx, cw = col_xs_t[ci], col_ws_t[ci]
    add_text_box(slide, h, cx + 0.05, 1.91, cw, 0.3,
                 font_size=12, bold=True, color=YELLOW)
for ri, (term, defn, ex) in enumerate(terms):
    bg = RGBColor(0x10, 0x25, 0x40) if ri % 2 == 0 else RGBColor(0x14, 0x2D, 0x50)
    y = 2.22 + ri * 0.34
    for cx, cw, txt in zip(col_xs_t, col_ws_t, [term, defn, ex]):
        add_rect(slide, cx, y, cw, 0.34, bg)
        add_text_box(slide, "  " + txt, cx + 0.05, y + 0.03,
                     cw - 0.1, 0.28, font_size=11, color=WHITE)

# How to identify FDs
add_rect(slide, 0.4, 3.97, 12.5, 0.32, RGBColor(0x0D, 0x20, 0x38))
add_text_box(slide, "How to Identify Functional Dependencies (5 Steps):",
             0.55, 3.99, 8.0, 0.28, font_size=12, bold=True, color=ACCENT_BLUE)
steps_text = [
    ("  Step 1: List all attributes in the table", 12, False, WHITE),
    ("  Step 2: Identify the primary key (single or composite)", 12, False, WHITE),
    ("  Step 3: For each non-key attribute — ask 'What uniquely determines this?'", 12, False, WHITE),
    ("  Step 4: Write down all discovered FDs", 12, False, WHITE),
    ("  Step 5: Check for partial FDs (→ 2NF issue) and transitive FDs (→ 3NF issue)", 12, False, YELLOW),
]
add_multiline(slide, steps_text, 0.4, 4.32, 12.5, 1.7)

slide.notes_slide.notes_text_frame.text = (
    "Stress that a determinant can be ANY attribute — not just the primary key. "
    "A non-key attribute determining another non-key attribute is the cause of "
    "transitive dependencies (3NF violation). Walk students through the 5 steps "
    "using a concrete example table."
)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 8 — NORMALIZATION PROCESS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "D. Normalization Process")
slide_title_bar(slide, "The Normalization Process — Step-by-Step Framework")

# Left: Flow diagram
flow = [
    ("RAW / UNNORMALIZED DATA (UNF)", ACCENT_BLUE),
    ("Step 1 → 1NF  :  Make all values atomic; remove repeating groups", MID_BLUE),
    ("Step 2 → 2NF  :  Remove partial dependencies (split composite-key tables)", MID_BLUE),
    ("Step 3 → 3NF  :  Remove transitive dependencies (split non-key → non-key FDs)", MID_BLUE),
    ("NORMALIZED DATABASE ✅", GREEN),
]
box_left_f = 0.4
box_w_f    = 7.5
box_h_f    = 0.58
start_y_f  = 1.32
for i, (label, color) in enumerate(flow):
    y = start_y_f + i * (box_h_f + 0.06)
    add_rect(slide, box_left_f, y, box_w_f, box_h_f, color)
    bold = (i == 0 or i == len(flow) - 1)
    add_text_box(slide, "   " + label, box_left_f + 0.08, y + 0.08,
                 box_w_f - 0.12, box_h_f - 0.1,
                 font_size=12, bold=bold, color=WHITE)
    if i < len(flow) - 1:
        add_text_box(slide, "▼", box_left_f + box_w_f / 2 - 0.2, y + box_h_f - 0.04,
                     0.4, 0.26, font_size=11, color=LIGHT_GRAY,
                     align=PP_ALIGN.CENTER)

# Right: Decomposition principles
add_rect(slide, 8.2, 1.32, 4.7, 0.36, RGBColor(0x0D, 0x30, 0x1A))
add_text_box(slide, "✅ Lossless-Join Decomposition",
             8.3, 1.34, 4.5, 0.32, font_size=13, bold=True, color=GREEN)
lossless = [
    "Joining all decomposed tables returns",
    "EXACTLY the original data — no extra or",
    "missing rows (no spurious tuples).",
    "  Test: Common attribute must be a key",
    "  in at least one of the split tables.",
]
add_multiline(slide, [(t, 11, False, WHITE) for t in lossless],
              8.2, 1.72, 4.7, 1.3)

add_rect(slide, 8.2, 3.08, 4.7, 0.36, RGBColor(0x0D, 0x20, 0x38))
add_text_box(slide, "✅ Dependency Preservation",
             8.3, 3.10, 4.5, 0.32, font_size=13, bold=True, color=ACCENT_BLUE)
dep_p = [
    "All original FDs must still be enforceable",
    "using the decomposed tables alone.",
    "If a dependency is lost, the database",
    "cannot enforce that business rule without",
    "additional application-level checks.",
]
add_multiline(slide, [(t, 11, False, WHITE) for t in dep_p],
              8.2, 3.48, 4.7, 1.4)

add_rect(slide, 8.2, 4.95, 4.7, 0.7, RGBColor(0x14, 0x2D, 0x50))
add_text_box(slide,
             "  Key Principles:\n"
             "  • Always work bottom-up (1NF → 2NF → 3NF)\n"
             "  • Lossless join is NON-NEGOTIABLE\n"
             "  • 3NF is the standard target for most systems",
             8.3, 4.97, 4.5, 0.66, font_size=10, color=LIGHT_GRAY)

slide.notes_slide.notes_text_frame.text = (
    "Present this as a practical workflow. Normalization is iterative — check each step "
    "before moving to the next. Lossless join is non-negotiable: if decomposition creates "
    "extra rows when tables are rejoined, the decomposition is wrong. Most real systems "
    "target 3NF as their baseline."
)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 9 — NORMAL FORMS COMPARISON TABLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "Summary")
slide_title_bar(slide, "Normal Forms at a Glance — Comparison Table")

# Comparison table: 1NF, 2NF, 3NF
summary_data = [
    ("1NF", "Atomic values, no repeating groups",
     "Multi-valued / repeating columns",
     "Single or composite PK with atomic values"),
    ("2NF", "Full dependency on entire primary key",
     "Partial dependencies (non-key attr depends on part of composite PK)",
     "Composite key: split partial FDs into separate tables"),
    ("3NF", "No non-key attribute depends on another non-key attribute",
     "Transitive dependencies (non-key → non-key → value)",
     "Move transitively dependent attrs to their own table"),
]
col_headers = ["NF", "Key Rule", "Removes / Fixes", "How to Achieve"]
col_xs = [0.4, 1.55, 4.65, 8.75]
col_ws = [1.1, 3.05, 4.05, 4.25]

add_rect(slide, 0.4, 1.32, 12.9, 0.38, MID_BLUE)
for cx, cw, h in zip(col_xs, col_ws, col_headers):
    add_text_box(slide, h, cx + 0.05, 1.33, cw - 0.1, 0.35,
                 font_size=12, bold=True, color=YELLOW)

nf_colors_s = [RGBColor(0xE7, 0x4C, 0x3C), RGBColor(0xF3, 0x9C, 0x12), RGBColor(0x27, 0xAE, 0x60)]
for ri, (row, nfc) in enumerate(zip(summary_data, nf_colors_s)):
    bg = RGBColor(0x10, 0x25, 0x40) if ri % 2 == 0 else RGBColor(0x14, 0x2D, 0x50)
    y = 1.70 + ri * 1.1
    for ci, (cx, cw, cell) in enumerate(zip(col_xs, col_ws, row)):
        add_rect(slide, cx, y, cw, 1.05, bg)
        col = nfc if ci == 0 else WHITE
        bold = ci == 0
        fs = 18 if ci == 0 else 11
        add_text_box(slide, "  " + cell, cx + 0.05, y + 0.05,
                     cw - 0.1, 0.95, font_size=fs, bold=bold, color=col)

add_rect(slide, 0.4, 5.12, 12.9, 0.52, RGBColor(0x14, 0x2D, 0x50))
add_text_box(slide,
             "  💡  Remember: Each higher NF builds on and includes all requirements of the forms below it.\n"
             "  A table in 3NF is also in 2NF and 1NF.",
             0.55, 5.14, 12.7, 0.48, font_size=12, bold=True, color=ACCENT_BLUE)

slide.notes_slide.notes_text_frame.text = (
    "Use this table as a quick reference. Students can work through any table by checking "
    "from 1NF upward. If a table fails 1NF, fix that first — you cannot skip levels. "
    "Emphasize that each NF is a superset of the previous one."
)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 10 — BEST PRACTICES & CONCLUSION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_chrome(slide, "Conclusion")
slide_title_bar(slide, "Best Practices & Key Takeaways")

# Key Takeaways
add_rect(slide, 0.4, 1.32, 12.5, 0.36, MID_BLUE)
add_text_box(slide, "🔑  Key Takeaways", 0.55, 1.34, 5.0, 0.32,
             font_size=14, bold=True, color=YELLOW)
takeaways = [
    "  •  1NF: Eliminate multi-valued columns — each cell holds exactly one value",
    "  •  2NF: Eliminate partial dependencies — every non-key attr depends on the full PK",
    "  •  3NF: Eliminate transitive dependencies — non-key attrs depend ONLY on the PK",
    "  •  Normalization is cumulative — a table in 3NF is also in 2NF and 1NF",
    "  •  Always verify lossless-join decomposition at every step",
]
add_multiline(slide, [(t, 13, False, WHITE) for t in takeaways], 0.4, 1.72, 12.5, 1.6)

# Best Practices vs Pitfalls
add_rect(slide, 0.4, 3.45, 12.5, 0.36, RGBColor(0x0D, 0x30, 0x1A))
add_text_box(slide, "✅  Best Practices", 0.5, 3.47, 6.0, 0.32,
             font_size=13, bold=True, color=GREEN)
add_text_box(slide, "❌  Common Pitfalls", 6.75, 3.47, 6.0, 0.32,
             font_size=13, bold=True, color=RED_SOFT)

bps = [
    "Start from UNF and normalize step by step",
    "Document all FDs before starting normalization",
    "Verify lossless join at every decomposition step",
    "Aim for 3NF as the baseline for production systems",
]
for i, bp in enumerate(bps):
    y = 3.85 + i * 0.34
    bg = RGBColor(0x0D, 0x30, 0x1A) if i % 2 == 0 else RGBColor(0x10, 0x38, 0x20)
    add_rect(slide, 0.4, y, 6.1, 0.32, bg)
    add_text_box(slide, "  •  " + bp, 0.45, y + 0.02, 6.0, 0.28, font_size=10, color=WHITE)

pitfalls = [
    ("Over-normalizing",        "Too many joins slow down queries"),
    ("Under-normalizing",       "Leads to update/insert/delete anomalies"),
    ("Ignoring composite keys", "Missing partial deps → 2NF error"),
    ("Breaking lossless join",  "Split tables cannot be rejoined correctly"),
]
for i, (mistake, reason) in enumerate(pitfalls):
    y = 3.85 + i * 0.34
    bg = RGBColor(0x30, 0x10, 0x10) if i % 2 == 0 else RGBColor(0x38, 0x14, 0x14)
    add_rect(slide, 6.75, y, 2.6, 0.32, bg)
    add_text_box(slide, "  " + mistake, 6.8, y + 0.02, 2.5, 0.28,
                 font_size=10, bold=True, color=YELLOW)
    add_rect(slide, 9.4, y, 3.5, 0.32, bg)
    add_text_box(slide, "  " + reason, 9.45, y + 0.02, 3.4, 0.28,
                 font_size=10, color=WHITE)

add_rect(slide, 0.4, 5.26, 12.5, 0.56, RGBColor(0x14, 0x2D, 0x50))
add_text_box(slide,
             "  💡  Normalization is a TOOL, not an end goal. Design with purpose!\n"
             "  \"Normalization is the foundation of good relational database design — "
             "every scalable, maintainable system relies on it.\"",
             0.55, 5.28, 12.2, 0.52, font_size=11, bold=True, color=ACCENT_BLUE)

slide.notes_slide.notes_text_frame.text = (
    "Wrap up by asking students to recall: what does each NF eliminate? "
    "Reiterate that normalization is cumulative. Over-normalization is a real-world "
    "problem — too many joins slow down read-heavy workloads. 3NF is the standard "
    "target for most OLTP systems."
)


# ═══════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════
output_file = "Advand_Normalization_Concepts.pptx"
prs.save(output_file)
print(f"✅  Presentation saved as:  {output_file}")
print(f"    Total slides: {len(prs.slides)}  (1 title + 9 content slides)")
